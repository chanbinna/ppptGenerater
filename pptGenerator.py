from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor

# Create a presentation object
prs = Presentation()

# Set slide dimensions for 25.4 cm width and 14.29 cm height
prs.slide_width = Cm(25.4)
prs.slide_height = Cm(14.29)

# Read the text file and parse lines
with open("text.txt", "r", encoding="utf-8") as file:
    lines = file.readlines()

# Initialize variables
current_title = ""
content_lines = []

# Iterate through each line
for line in lines:
    line = line.strip()  # Remove leading and trailing whitespace
    
    if line.startswith("<") and line.endswith(">"):
        # If a new title is found, create slides for the previous title and content
        if content_lines:
            # Process content lines in pairs to create slides
            for i in range(0, len(content_lines), 2):
                # Create a new slide
                slide = prs.slides.add_slide(prs.slide_layouts[6])

                # Add the background image
                img_path = "image/pptBackground.jpg"  # Replace with your image file path
                background = slide.shapes.add_picture(img_path, Cm(0), Cm(0), width=prs.slide_width, height=prs.slide_height)
                slide.shapes._spTree.remove(background._element)
                slide.shapes._spTree.insert(2, background._element)

                # Add top-left title
                top_left_box = slide.shapes.add_textbox(Cm(0.88), Cm(0.81), Cm(10), Cm(2))
                top_left_frame = top_left_box.text_frame
                top_left_frame.text = current_title
                top_left_paragraph = top_left_frame.paragraphs[0]
                top_left_paragraph.font.name = "Malgun Gothic"
                top_left_paragraph.font.bold = True
                top_left_paragraph.font.size = Pt(17.5)
                top_left_paragraph.font.color.rgb = RGBColor(0x59, 0x59, 0x59)
                top_left_paragraph.alignment = PP_ALIGN.LEFT

                # Add main centered content text box
                title_box = slide.shapes.add_textbox(Cm(0), Cm((14.29 - 2) / 2), prs.slide_width, Cm(2))
                title_frame = title_box.text_frame

                # Add each line of content as a separate paragraph with formatting
                for text_line in content_lines[i:i + 2]:
                    p = title_frame.add_paragraph() if title_frame.text else title_frame.paragraphs[0]
                    p.text = text_line
                    p.font.name = "Malgun Gothic"
                    p.font.size = Pt(40)
                    p.font.bold = True
                    p.alignment = PP_ALIGN.CENTER
                    p.font.color.rgb = RGBColor(0x1F, 0x38, 0x64)
                title_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

            # Reset content lines after slides are created for the current title
            content_lines = []
        
        # Update the current title to the new title found
        current_title = line[1:-1]  # Remove the angle brackets

    else:
        # Collect content lines under the current title
        if line:
            content_lines.append(line)

# Create slides for any remaining content after the last title
if content_lines:
    for i in range(0, len(content_lines), 2):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Add background image
        img_path = "image/pptBackground.jpg"
        background = slide.shapes.add_picture(img_path, Cm(0), Cm(0), width=prs.slide_width, height=prs.slide_height)
        slide.shapes._spTree.remove(background._element)
        slide.shapes._spTree.insert(2, background._element)
        
        # Add top-left title
        top_left_box = slide.shapes.add_textbox(Cm(0.88), Cm(0.81), Cm(10), Cm(2))
        top_left_frame = top_left_box.text_frame
        top_left_frame.text = current_title
        top_left_paragraph = top_left_frame.paragraphs[0]
        top_left_paragraph.font.name = "Malgun Gothic"
        top_left_paragraph.font.bold = True
        top_left_paragraph.font.size = Pt(17.5)
        top_left_paragraph.font.color.rgb = RGBColor(0x59, 0x59, 0x59)
        top_left_paragraph.alignment = PP_ALIGN.LEFT

        # Add main centered content text box
        title_box = slide.shapes.add_textbox(Cm(0), Cm((14.29 - 2) / 2), prs.slide_width, Cm(2))
        title_frame = title_box.text_frame

        # Add each line of content as a separate paragraph with formatting
        for text_line in content_lines[i:i + 2]:
            p = title_frame.add_paragraph() if title_frame.text else title_frame.paragraphs[0]
            p.text = text_line
            p.font.name = "Malgun Gothic"
            p.font.size = Pt(40)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            p.font.color.rgb = RGBColor(0x1F, 0x38, 0x64)
        title_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

# Save the presentation
prs.save("찬양.pptx")